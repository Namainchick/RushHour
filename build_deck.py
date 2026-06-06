#!/usr/bin/env python3
"""Build the RushHour pitch deck (7 slides, 16:9) from the pitch script.

Run: .pptx-venv/bin/python build_deck.py
Output: RushHour_Pitch.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Brand ---
BG = RGBColor(0x0E, 0x11, 0x16)        # near-black
INK = RGBColor(0xFF, 0xFF, 0xFF)       # white
MUTE = RGBColor(0x9A, 0xA4, 0xB2)      # muted grey
ACCENT = RGBColor(0xFF, 0x5A, 0x36)    # rush orange-red
ACCENT2 = RGBColor(0x4C, 0x8B, 0xFF)   # blue
FONT = "Arial"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, EMU_W, EMU_H)  # rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # push to back
    sp = bg._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def txt(s, left, top, width, height, text, size, color=INK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, spacing=None):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.name = font
        f.color.rgb = color
    return tb


def kicker(s, text):
    txt(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.5),
        text.upper(), 14, ACCENT, bold=True)


def brand(s, n):
    txt(s, Inches(11.6), Inches(6.95), Inches(1.6), Inches(0.4),
        f"RushHour · {n}/7", 10, MUTE, align=PP_ALIGN.RIGHT)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def chip(s, left, top, width, text, color, fill=None):
    box = s.shapes.add_shape(5, left, top, width, Inches(0.9))  # rounded rect
    box.fill.solid()
    box.fill.fore_color.rgb = fill if fill else BG
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = INK
    return box


# ============ SLIDE 1 — Hook ============
s = slide()
kicker(s, "RushHour")
txt(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
    "Today, anyone can build anything with AI.", 40, INK, bold=True, spacing=1.0)
txt(s, Inches(0.7), Inches(2.7), Inches(5.7), Inches(2.0),
    "Building is easy now.\n→  Distribution wins.", 26, ACCENT2, bold=True, spacing=1.1)
txt(s, Inches(6.9), Inches(2.7), Inches(5.7), Inches(2.0),
    "AI makes everything generic.\n→  Real humans win.", 26, ACCENT2, bold=True, spacing=1.1)
txt(s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.5),
    "= Social Media Marketing", 44, ACCENT, bold=True, align=PP_ALIGN.CENTER)
brand(s, 1)
notes(s,
"Today, anyone can build anything with AI. Look around this room - in eight hours, "
"teams here built incredible things. Building used to be the hard part. It isn't "
"anymore. So the product matters less than ever - what matters now is distribution: "
"reaching the right people.\n\n"
"But there's a second shift. Everything AI makes looks the same. Generic. Soulless "
"ads nobody remembers. What actually cuts through today is the opposite - real, "
"organic content from real humans.\n\n"
"Put those two together, and they point to one thing: social media marketing.")

# ============ SLIDE 2 — Problem ============
s = slide()
kicker(s, "The Problem")
txt(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
    "Social media marketing is broken.", 36, INK, bold=True)
chip(s, Inches(0.7), Inches(2.9), Inches(3.6),
     "BUSINESSES\ncan't find the right creators\nslow + expensive", ACCENT2)
chip(s, Inches(4.85), Inches(2.7), Inches(3.6),
     "AGENCY\nhuge cut just to introduce them", ACCENT, fill=RGBColor(0x2A,0x12,0x0C))
chip(s, Inches(9.0), Inches(2.9), Inches(3.6),
     "CREATORS\nwant gigs, can't find clients", ACCENT2)
txt(s, Inches(0.7), Inches(4.6), Inches(12), Inches(0.6),
    "Business  →  💸 Agency  →  Creator", 22, MUTE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.0),
    "Hours of hunting, big budgets, and a middleman taking a cut.", 20, INK,
    align=PP_ALIGN.CENTER)
brand(s, 2)
notes(s,
"But social media marketing is broken. On one side, businesses spend hours and real "
"budget hunting for creators that actually fit them. On the other, thousands of "
"creators want gigs - and can't find clients. And right in the middle sits an agency, "
"taking a huge cut just to introduce them.")

# ============ SLIDE 3 — Solution ============
s = slide()
kicker(s, "The Solution")
txt(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
    "We cut the agency out — and replace it with AI.", 34, INK, bold=True)
txt(s, Inches(0.7), Inches(2.5), Inches(12), Inches(1.2),
    "Airbnb for social media marketing.", 40, ACCENT, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(4.0), Inches(12), Inches(0.7),
    "Business  →  [ RushHour AI ]  →  Creator", 24, INK, bold=True, align=PP_ALIGN.CENTER)
chip(s, Inches(1.3), Inches(5.1), Inches(3.2), "AI matchmaking", ACCENT2)
chip(s, Inches(5.05), Inches(5.1), Inches(3.2), "Payments", ACCENT2)
chip(s, Inches(8.8), Inches(5.1), Inches(3.2), "Trust + reviews", ACCENT2)
brand(s, 3)
notes(s,
"So we cut the agency out - and replace it with AI. Think of us as Airbnb for social "
"media marketing. Airbnb killed the travel agency by matching people directly and "
"taking a small cut. We do the same for brands and creators - except our matchmaker "
"is an AI agent. And everything Airbnb has - payments, trust, reviews - we have too.")

# ============ SLIDE 4 — Demo ============
s = slide()
kicker(s, "Live Demo")
txt(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
    "From brand to best-fit creators — in seconds.", 32, INK, bold=True)
chip(s, Inches(0.7), Inches(2.7), Inches(2.7), "1 · Drop in website + brand", ACCENT2)
chip(s, Inches(3.6), Inches(2.7), Inches(2.7), "2 · AI builds brand profile", ACCENT2)
chip(s, Inches(6.5), Inches(2.7), Inches(2.7), "3 · Scores every creator", ACCENT2)
chip(s, Inches(9.4), Inches(2.7), Inches(2.7), "4 · Top 3 matches", ACCENT, fill=RGBColor(0x2A,0x12,0x0C))
txt(s, Inches(0.7), Inches(4.4), Inches(12), Inches(1.0),
    "▶  LIVE DEMO", 40, ACCENT, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1.0),
    "Not a chatbot — AI agents analyze brand, audience, content style & past collabs.\nThe AI is the product.",
    20, MUTE, align=PP_ALIGN.CENTER, spacing=1.1)
brand(s, 4)
notes(s,
"Let me show you. A restaurant drops in their website and brand. Our AI builds a "
"structured profile in seconds - then scores every creator for fit. (click) And "
"here's the result: the top three creators most likely to actually work for them. "
"No scrolling, no guessing.\n\n"
"And under the hood, this isn't a chatbot - it's AI agents analyzing brand, audience, "
"content style, and past collaborations. The AI is the product.")

# ============ SLIDE 5 — Market ============
s = slide()
kicker(s, "Market")
txt(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
    "Is this a real market? Yes.", 36, INK, bold=True)
txt(s, Inches(0.7), Inches(2.9), Inches(6.0), Inches(1.5),
    "$[ __ ]B", 60, ACCENT, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(4.3), Inches(6.0), Inches(0.7),
    "influencer marketing, growing yearly", 18, MUTE, align=PP_ALIGN.CENTER)
txt(s, Inches(6.9), Inches(2.9), Inches(6.0), Inches(1.5),
    "[ __ ]%", 60, ACCENT, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(6.9), Inches(4.3), Inches(6.0), Inches(0.7),
    "of brands already use creators", 18, MUTE, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.8),
    "But honestly? Forget the theory.", 28, INK, bold=True, align=PP_ALIGN.CENTER)
brand(s, 5)
notes(s,
"Is this a real market? Yes - influencer marketing is a [___]-billion-dollar industry, "
"growing every year. But honestly? Forget the theory.\n\n"
"[Fill in real numbers from Jusuf/Kobi before submission.]")

# ============ SLIDE 6 — Scouting ============
s = slide()
kicker(s, "Our Real Data — Scouting")
txt(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
    "We went out into the real world.", 36, INK, bold=True)
txt(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.6),
    "Not surveys. Actual conversations. Both sides.", 20, MUTE)
chip(s, Inches(0.7), Inches(3.2), Inches(5.8),
     "BUSINESSES\n[ __ ] restaurants scouted", ACCENT2)
chip(s, Inches(6.8), Inches(3.2), Inches(5.8),
     "CREATORS\ninfluencers on board — incl. big ones (2 clips)", ACCENT2)
txt(s, Inches(0.7), Inches(4.6), Inches(12), Inches(1.0),
    "FIRST DEALS CLOSED ✅", 40, ACCENT, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.8),
    "We didn't ask “would you like this?” — we asked for the deal. They're paying.",
    20, INK, align=PP_ALIGN.CENTER)
brand(s, 6)
notes(s,
"We went out into the real world - not surveys, actual conversations. We scouted both "
"sides. On the creator side, we already have influencers on board - including big "
"ones. (show clips) On the business side, we didn't ask 'would you like this?' - we "
"asked for the deal. And we already closed our first ones. People don't just want "
"this. They're paying.\n\n"
"[Fill in real numbers: X restaurants approached, Y deals closed, Z creators.]")

# ============ SLIDE 7 — Niche + Why us + Close ============
s = slide()
kicker(s, "Why Us")
txt(s, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
    "We start where distribution decides survival:", 30, INK, bold=True)
txt(s, Inches(0.7), Inches(2.2), Inches(12), Inches(1.0),
    "Restaurants & Startups.", 40, ACCENT, bold=True)
txt(s, Inches(0.7), Inches(3.6), Inches(12), Inches(1.2),
    "4 founders — all creators. We run our own social media.\nWe live this problem, on both sides.",
    24, INK, bold=True, spacing=1.1)
txt(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.4),
    "Anyone can build. The hard part is knowing if anyone wants it.\nWe went out, we asked — and they said yes. That's RushHour.",
    24, ACCENT, bold=True, align=PP_ALIGN.CENTER, spacing=1.1)
brand(s, 7)
notes(s,
"We start where distribution decides survival: restaurants and startups. And why us? "
"We're four founders who are creators ourselves. We run our own social media - we know "
"exactly where this hurts, on both sides, because we live it.\n\n"
"Today, anyone can build. The hard part is knowing if anyone wants it. We went out, "
"we asked - and they said yes. That's RushHour.")

prs.save("RushHour_Pitch.pptx")
print("Saved RushHour_Pitch.pptx with", len(prs.slides._sldIdLst), "slides")
